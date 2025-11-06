import cv2
import easyocr
import numpy as np
from googletrans import Translator
from PIL import Image, ImageDraw, ImageFont
import base64
import os
from io import BytesIO
from collections import OrderedDict

from flask import Flask, request, jsonify, send_from_directory, render_template_string, send_file
from flask_cors import CORS

import docx
import pptx
import pdfplumber
from fpdf import FPDF
from deep_translator import GoogleTranslator as DeepTranslator 

app = Flask(__name__)
CORS(app)


FONT_PATH = "NotoSans-Regular.ttf" 


reader = easyocr.Reader(['en'], gpu=False) 
text_translator = Translator()

MAX_CACHE_SIZE = 200
translation_cache = OrderedDict()


def translate_text_image(text, dest_language='vi'):
    """
    Dịch văn bản sang ngôn ngữ được chỉ định bằng Google Translate API.
    """
    try:
        translator = Translator()
        translation = translator.translate(text, dest=dest_language)
        return translation.text
    except Exception as e:
        print(f"Lỗi dịch: {e}")
        return text  # Trả về văn bản gốc nếu dịch lỗi

def blur_bbox(img, bbox):
  
    mask = np.zeros(img.shape[:2], dtype=np.uint8)
    cv2.fillPoly(mask, [bbox], 255)
    blurred_img = cv2.GaussianBlur(img, (51, 51), 0)
    img_blurred_area = np.where(mask[:, :, np.newaxis] == 255, blurred_img, img)
    return img_blurred_area

def draw_text_utf8(image, text, position, font_path, font_size, color=(255, 0, 0)):
  
    pil_img = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
    draw = ImageDraw.Draw(pil_img)

    try:
        font = ImageFont.truetype(font_path, font_size)
    except IOError:
        print(f"Lỗi: Không tìm thấy tệp phông chữ tại '{font_path}'. Đảm bảo bạn đã tải xuống và đặt phông chữ hỗ trợ Unicode vào đúng thư mục. Sử dụng phông chữ mặc định.")
        try:
            font = ImageFont.truetype("arial.ttf", font_size) # Thử Arial nếu Noto Sans thất bại
        except IOError:
            font = ImageFont.load_default() # Phương án cuối cùng
            print("Lỗi: Không tìm thấy phông chữ Arial. Sử dụng phông chữ mặc định của Pillow.")

    draw.text(position, text, font=font, fill=color)
    return cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)


def translate_text_batch(text_list, target_lang):

    if not text_list:
        return []
    try:
        separator = "\n<br/>\n"
        combined_text = separator.join(text_list)
        translator = DeepTranslator(source='auto', target=target_lang)
        translated_combined = translator.translate(combined_text)
        if separator in translated_combined:
            return translated_combined.split(separator)
        else:
            return [DeepTranslator(source='auto', target=target_lang).translate(text) for text in text_list]
    except Exception as e:
        print(f"Lỗi khi dịch hàng loạt: {e}")
        return [DeepTranslator(source='auto', target=target_lang).translate(text) for text in text_list]

def handle_txt(file_stream, target_lang):
    """Dịch file .txt"""
    content = file_stream.read().decode('utf-8', errors='ignore')
    translated_content = DeepTranslator(source='auto', target=target_lang).translate(content)
    output_stream = BytesIO()
    output_stream.write(translated_content.encode('utf-8'))
    output_stream.seek(0)
    return output_stream

def handle_docx(file_stream, target_lang):
    """Dịch file .docx và giữ nguyên định dạng."""
    doc = docx.Document(file_stream)

    texts_to_translate = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts_to_translate.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts_to_translate.append(cell.text)

    translated_texts = translate_text_batch(texts_to_translate, target_lang)

    translated_iter = iter(translated_texts)
    for para in doc.paragraphs:
        if para.text.strip():
            try:
                translated_text = next(translated_iter)
                para.text = ''
                para.add_run(translated_text)
            except StopIteration:
                break
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    try:
                        cell.text = next(translated_iter)
                    except StopIteration:
                        break

    output_stream = BytesIO()
    doc.save(output_stream)
    output_stream.seek(0)
    return output_stream

def handle_pptx(file_stream, target_lang):

    prs = pptx.Presentation(file_stream)
    texts_to_translate = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    texts_to_translate.append(para.text)

    translated_texts = translate_text_batch(texts_to_translate, target_lang)
    translated_iter = iter(translated_texts)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    try:
                        para.text = next(translated_iter)
                    except StopIteration:
                        break

    output_stream = BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    return output_stream


HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="vi">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Ứng dụng dịch thuật đa năng</title>
    <link href="https://fonts.googleapis.com/css2?family=Roboto:wght@400;700&family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
    <style>
        :root {
            --primary-color: #5b6ee1;
            --secondary-color: #38c6e5;
            --primary-gradient: linear-gradient(100deg, #5b6ee1 0%, #38c6e5 100%);
            --primary-gradient-hover: linear-gradient(100deg, #38c6e5 0%, #5b6ee1 100%);
            --background-color: #f7faff;
            --surface-color: #ffffff;
            --text-color: #232946;
            --text-color-light: #6b7280;
            --border-color: #e3e8f7;
            --border-radius: 18px;
            --shadow: 0 8px 32px 0 rgba(91,110,225,0.10), 0 2px 8px 0 rgba(56,198,229,0.08);
            --focus-shadow: 0 0 0 4px rgba(91,110,225,0.13);
        }

        body {
            font-family: 'Inter', 'Roboto', Arial, sans-serif;
            background: var(--background-color);
            color: var(--text-color);
            margin: 0;
            padding: 24px;
            min-height: 100vh;
            display: flex;
            align-items: flex-start; /* Changed to flex-start to allow scrolling */
            justify-content: center;
            box-sizing: border-box;
            animation: fadeInBg 1.2s;
            overflow-y: auto; /* Enable vertical scrolling */
        }

        @keyframes fadeInBg {
            from { background: #e3e8f7; }
            to { background: var(--background-color);}
        }

        .container {
            max-width: 900px; /* Increased max-width for better layout */
            width: 100%;
            background: var(--surface-color);
            padding: 40px 34px 34px 34px;
            border-radius: var(--border-radius);
            box-shadow: var(--shadow);
            position: relative;
            overflow: hidden;
            animation: popIn 0.8s cubic-bezier(.68,-0.55,.27,1.55);
            border: 2.5px solid #eaf6ff;
            margin-top: 20px; /* Add some top margin */
            margin-bottom: 20px; /* Add some bottom margin */
        }

        @keyframes popIn {
            0% { transform: scale(0.93) translateY(30px); opacity: 0;}
            100% { transform: scale(1) translateY(0); opacity: 1;}
        }

        h1 {
            text-align: center;
            margin-top: 0;
            margin-bottom: 32px;
            font-weight: 900;
            font-size: 2.5rem;
            letter-spacing: 1px;
            background: var(--primary-gradient);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
            animation: gradientMove 2.5s linear infinite alternate;
        }

        @keyframes gradientMove {
            0% { background-position: 0% 50%;}
            100% { background-position: 100% 50%;}
        }

        .tab-container {
            display: flex;
            justify-content: center;
            margin-bottom: 30px;
            background: #f0f4f8;
            border-radius: 12px;
            padding: 8px;
            box-shadow: inset 0 1px 4px rgba(0,0,0,0.05);
        }

        .tab-button {
            flex: 1;
            padding: 12px 20px;
            border: none;
            background: transparent;
            font-size: 1.1em;
            font-weight: 600;
            color: var(--text-color-light);
            cursor: pointer;
            border-radius: 10px;
            transition: all 0.3s ease;
            text-align: center;
        }

        .tab-button.active {
            background: var(--primary-gradient);
            color: #fff;
            box-shadow: 0 4px 15px rgba(91,110,225,0.2);
            transform: translateY(-2px);
        }

        .tab-button:hover:not(.active) {
            background: #e6eaf0;
            color: var(--primary-color);
        }

        .tab-content {
            display: none;
            padding: 20px 0;
        }

        .tab-content.active {
            display: block;
        }

        /* Styles for Text Translator */
        .lang-selectors {
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 20px;
            gap: 16px;
            background: linear-gradient(90deg, #eaf6ff 0%, #f7faff 100%);
            border-radius: 14px;
            padding: 12px 16px;
            box-shadow: 0 2px 12px rgba(91,110,225,0.07);
            animation: fadeIn 1.2s;
        }

        .lang-selectors select {
            flex: 1;
            cursor: pointer;
            background: #fff;
            border: 1.5px solid var(--border-color);
            border-radius: 10px;
            font-size: 17px;
            padding: 12px 38px 12px 16px;
            color: var(--text-color);
            transition: border-color 0.2s, box-shadow 0.2s, background 0.3s;
            background-image: url("data:image/svg+xml,%3csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 16 16'%3e%3cpath fill='none' stroke='%235b6ee1' stroke-linecap='round' stroke-linejoin='round' stroke-width='2' d='M2 5l6 6 6-6'/%3e%3c/svg%3e");
            background-repeat: no-repeat;
            background-position: right 1rem center;
            background-size: 18px 14px;
            box-shadow: 0 1px 6px rgba(91,110,225,0.04);
            appearance: none;
            -webkit-appearance: none;
            -moz-appearance: none;
        }
        .lang-selectors select:focus {
            border-color: var(--primary-color);
            box-shadow: var(--focus-shadow);
            outline: none;
            background: #eaf6ff;
        }
        .lang-selectors select:hover {
            background: #f0fcfa;
        }

        .lang-selectors span {
            font-size: 26px;
            font-weight: bold;
            color: var(--primary-color);
            transition: transform 0.25s cubic-bezier(.68,-0.55,.27,1.55), color 0.2s;
            user-select: none;
            filter: drop-shadow(0 2px 4px #38c6e533);
        }
        .lang-selectors span:hover {
            transform: scale(1.25) rotate(12deg);
            color: var(--secondary-color);
        }

        .input-group {
            display: flex;
            gap: 16px;
            margin-bottom: 20px;
            position: relative;
        }

        textarea {
            width: 100%;
            padding: 18px;
            font-size: 18px;
            font-family: 'Inter', sans-serif;
            border-radius: 12px;
            border: 1.5px solid var(--border-color);
            background: #fafdff;
            transition: border-color 0.2s, box-shadow 0.2s, background 0.3s;
            resize: vertical;
            min-height: 120px;
            box-shadow: 0 1px 8px rgba(91,110,225,0.07);
            outline: none;
        }
        textarea:focus {
            border-color: var(--primary-color);
            box-shadow: 0 0 0 4px #5b6ee122;
            background: #eaf6ff;
        }
        textarea:hover {
            background: #f0fcfa;
        }

        .mic-btn {
            background: var(--primary-gradient);
            border: none;
            color: #fff;
            border-radius: 50%;
            width: 56px; height: 56px;
            min-width: 56px;
            display: flex; align-items: center; justify-content: center;
            font-size: 28px;
            cursor: pointer;
            box-shadow: 0 4px 18px rgba(56,198,229,0.13), 0 1.5px 6px rgba(91,110,225,0.10);
            transition: background 0.25s, transform 0.18s, box-shadow 0.22s;
            align-self: flex-start;
            position: relative;
            z-index: 2;
            outline: none;
            animation: micAppear 1.2s;
        }
        @keyframes micAppear {
            from { opacity: 0; transform: scale(0.7);}
            to { opacity: 1; transform: scale(1);}
        }
        .mic-btn:hover {
            background: var(--primary-gradient-hover);
            transform: scale(1.12) rotate(-6deg);
            box-shadow: 0 8px 32px rgba(56,198,229,0.18);
        }
        .mic-btn:active {
            transform: scale(0.96);
        }
        .mic-btn.listening {
            background: linear-gradient(135deg, #ff7f9c 0%, #ffe347 100%);
            color: #fff;
            animation: pulse 1.2s infinite;
            box-shadow: 0 0 0 0 rgba(255, 127, 156, 0.5);
        }
        @keyframes pulse {
            0% { box-shadow: 0 0 0 0 rgba(255, 127, 156, 0.5);}
            70% { box-shadow: 0 0 0 18px rgba(255, 127, 156, 0);}
            100% { box-shadow: 0 0 0 0 rgba(255, 127, 156, 0);}
        }

        button[type="submit"], .action-button {
            background: var(--primary-gradient);
            color: #fff;
            border: none;
            padding: 16px 0;
            border-radius: 14px;
            font-size: 20px;
            font-weight: 900;
            width: 100%;
            cursor: pointer;
            margin-top: 12px;
            box-shadow: 0 2px 14px rgba(91,110,225,0.10);
            transition: background 0.22s, transform 0.13s, box-shadow 0.22s;
            letter-spacing: 0.7px;
            position: relative;
            overflow: hidden;
        }
        button[type="submit"]:hover, .action-button:hover {
            background: var(--primary-gradient-hover);
            transform: translateY(-2px) scale(1.04);
            box-shadow: 0 8px 28px rgba(56,198,229,0.13);
        }
        button[type="submit"]:active, .action-button:active {
            transform: scale(0.98);
        }
        button:disabled, .action-button:disabled {
            background: #cfd8dc;
            color: #888;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }

        .result {
            margin-top: 32px;
            padding: 24px 20px;
            background: linear-gradient(90deg, #fafdff 0%, #eaf6ff 100%);
            border-radius: 14px;
            font-size: 20px;
            min-height: 54px;
            color: var(--text-color);
            line-height: 1.7;
            word-wrap: break-word;
            box-shadow: 0 2px 16px rgba(91,110,225,0.08);
            border: 1.5px solid #e6eaff;
            animation: fadeIn 1.2s;
            transition: background 0.3s, box-shadow 0.3s;
            position: relative;
            overflow: hidden;
        }
        .result:empty {
            background: none;
            border: none;
            box-shadow: none;
            min-height: 0;
            padding: 0;
        }

        .result::after {
            content: "";
            display: block;
            position: absolute;
            left: -40px; top: -40px;
            width: 120px; height: 120px;
            background: radial-gradient(circle, #38c6e522 0%, transparent 80%);
            z-index: 0;
            pointer-events: none;
            animation: floatGlow 3s infinite alternate;
        }
        @keyframes floatGlow {
            0% { transform: translateY(0);}
            100% { transform: translateY(18px);}
        }

        .swap-btn {
            background: #fff;
            border: 1.5px solid var(--border-color);
            border-radius: 50%;
            width: 44px;
            height: 44px;
            font-size: 22px;
            color: var(--primary-color);
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            margin: 0 8px;
            transition:
                background 0.18s,
                color 0.18s,
                box-shadow 0.18s;
            box-shadow: 0 2px 8px rgba(91,110,225,0.09);
            outline: none;
            position: relative;
            z-index: 1;
            filter: drop-shadow(0 2px 8px #38c6e533);
        }
        .swap-btn:hover {
            background: var(--primary-gradient);
            color: #fff;
            box-shadow: 0 6px 18px rgba(56,198,229,0.18);
            filter: brightness(1.1) drop-shadow(0 4px 16px #5b6ee144);
        }
        .swap-btn:active {
            filter: brightness(1.15) drop-shadow(0 6px 18px #38c6e544);
        }
        .swap-btn .swap-icon {
            display: inline-block;
            transition: transform 0.5s cubic-bezier(.68,-0.55,.27,1.55);
        }
        .swap-btn.swapping .swap-icon {
            transform: rotate(360deg);
        }
        .lang-selectors select.swapping {
            border-color: var(--primary-color);
            background: #eaf6ff;
            box-shadow: 0 0 0 4px #5b6ee122, 0 2px 12px #38c6e533;
            transform: scale(1.07);
            transition:
                border-color 0.25s,
                background 0.25s,
                box-shadow 0.25s,
                transform 0.25s;
        }
        .lang-selectors select.swapped {
            border-color: var(--secondary-color);
            background: #f0fcfa;
            box-shadow: 0 0 0 2px #38c6e522;
            transform: scale(1.03);
            transition:
                border-color 0.25s,
                background 0.25s,
                box-shadow 0.25s,
                transform 0.25s;
        }

        /* Styles for Image Translator */
        .image-translator-section .input-group {
            flex-direction: column;
            align-items: center;
        }
        .image-translator-section input[type="file"] {
            display: block;
            margin: 0 auto 18px auto;
            padding: 10px;
            border: 1px solid var(--border-color);
            border-radius: 7px;
            background: #fff;
            font-size: 1em;
            width: 80%; /* Adjusted width */
            max-width: 400px; /* Max width for file input */
        }
        .image-translator-section .action-button {
            margin-bottom: 15px;
            width: 80%;
            max-width: 400px;
        }
        .image-translator-section hr {
            border: 0;
            height: 1px;
            background: #e3eaf1;
            margin: 28px 0;
        }
        .image-translator-section #loading {
            margin-top: 18px;
            font-style: italic;
            color: var(--primary-color);
            font-size: 1.13em;
            animation: pulse 1.5s infinite;
            text-align: center;
        }
        @keyframes pulse {
            0% { opacity: 1; }
            50% { opacity: 0.5; }
            100% { opacity: 1; }
        }
        .image-display-section {
            display: flex;
            flex-wrap: wrap;
            justify-content: center;
            gap: 28px;
            margin-top: 28px;
        }
        .image-container {
            flex: 1;
            min-width: 280px;
            max-width: 420px;
            border: 1px solid var(--border-color);
            border-radius: 12px;
            padding: 18px;
            background: #f7fafd;
            text-align: center;
            box-shadow: 0 2px 8px rgba(91,134,229,0.05);
        }
        .image-container h2 {
            color: var(--secondary-color);
            margin-top: 0;
            margin-bottom: 12px;
            font-size: 1.18em;
            background: none; /* Override global h2 gradient */
            -webkit-background-clip: unset;
            -webkit-text-fill-color: unset;
            background-clip: unset;
            animation: none;
        }
        .image-container img, .image-container video {
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            margin-top: 8px;
            box-shadow: 0 2px 8px rgba(91,134,229,0.08);
            display: block;
            margin-left: auto;
            margin-right: auto;
        }
        #translatedTextResults {
            margin-top: 32px;
            padding: 22px;
            border: 1px solid var(--border-color);
            border-radius: 12px;
            background: #f7fafd;
            text-align: left;
            box-shadow: 0 2px 8px rgba(91,134,229,0.05);
        }
        #translatedTextResults h2 {
            text-align: center;
            color: var(--text-color);
            margin-bottom: 16px;
            font-size: 1.25em;
            background: none; /* Override global h2 gradient */
            -webkit-background-clip: unset;
            -webkit-text-fill-color: unset;
            background-clip: unset;
            animation: none;
        }
        #resultsList {
            list-style: none;
            padding: 0;
        }
        #resultsList li {
            background: #e0f7fa;
            margin-bottom: 10px;
            padding: 12px 18px;
            border-radius: 7px;
            border-left: 6px solid var(--secondary-color);
            font-size: 1em;
            color: var(--text-color);
            box-shadow: 0 1px 3px rgba(91,134,229,0.04);
        }
        .hidden {
            display: none !important; /* Use !important to ensure it overrides other styles */
        }
        #webcamFeed {
            display: block;
            margin: 0 auto 12px auto;
            border-radius: 12px;
            box-shadow: 0 4px 24px rgba(91,134,229,0.18);
            border: 2px solid var(--secondary-color);
            max-width: 100%;
            background: #e0eafc;
            transition: box-shadow 0.3s;
        }
        #captureButton, #startWebcamButton {
            margin: 0 8px 12px 0;
        }
        .image-translator-section .input-group label {
            font-weight: 500;
            color: var(--text-color);
            font-size: 1.07em;
        }
        .image-translator-section select, .image-translator-section input[type="number"] {
            padding: 8px 14px;
            border: 1px solid var(--border-color);
            border-radius: 6px;
            font-size: 1em;
            background: #f7fafd;
            min-width: 120px;
        }

        /* Styles for File Translator */
        .file-translator-section .form-group {
            margin-bottom: 1.5rem;
            text-align: left;
        }
        .file-translator-section label {
            display: block;
            margin-bottom: 0.5rem;
            font-weight: 600;
            color: var(--text-color);
        }
        .file-translator-section input[type="file"], .file-translator-section select {
            width: 100%;
            padding: 0.8rem;
            border: 1px solid var(--border-color);
            border-radius: 6px;
            box-sizing: border-box;
            font-size: 1rem;
            background: #fafdff;
        }
        .file-translator-section .action-button {
            margin-top: 1rem;
        }
        .file-translator-section #loading-file {
            margin-top: 2rem;
            text-align: center;
        }
        .file-translator-section .spinner {
            border: 5px solid rgba(0, 0, 0, 0.1);
            width: 40px;
            height: 40px;
            border-radius: 50%;
            border-left-color: var(--primary-color);
            animation: spin 1s ease infinite;
            margin: 0 auto 1rem;
        }
        @keyframes spin { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }
        .file-translator-section #status-file {
            font-weight: 500;
            color: var(--text-color-light);
        }
        .file-translator-section .file-types {
            font-size: 0.9rem;
            color: var(--text-color-light);
            margin-top: -1.5rem;
            margin-bottom: 2rem;
            text-align: center;
        }
        .file-translator-section #download-section {
            margin-top: 2rem;
            border-top: 1px solid #eee;
            padding-top: 1.5rem;
            text-align: center;
        }
        .file-translator-section #download-link, .file-translator-section #view-button {
            display: inline-block;
            background: var(--primary-color);
            color: white;
            padding: 0.8rem 1.5rem;
            border-radius: 6px;
            text-decoration: none;
            font-weight: bold;
            transition: background 0.3s;
            margin-right: 10px;
            margin-bottom: 10px; /* For mobile spacing */
        }
        .file-translator-section #download-link:hover, .file-translator-section #view-button:hover {
            background: var(--primary-gradient-hover);
        }
        .file-translator-section #translated-text-display {
            margin-top: 2rem;
            padding: 1rem;
            border: 1px dashed var(--border-color);
            border-radius: 8px;
            text-align: left;
            max-height: 200px;
            overflow-y: auto;
            background-color: #f9f9f9;
            white-space: pre-wrap;
            font-size: 0.9rem;
            color: var(--text-color);
        }

        /* Responsive adjustments */
        @media (max-width: 768px) {
            .container {
                padding: 20px 15px;
                margin-top: 10px;
                margin-bottom: 10px;
            }
            h1 {
                font-size: 2rem;
            }
            .tab-button {
                padding: 10px 15px;
                font-size: 0.95em;
            }
            .lang-selectors {
                flex-direction: column;
                gap: 10px;
                padding: 10px;
            }
            .lang-selectors select {
                width: 100%;
                padding: 10px 30px 10px 12px;
                font-size: 16px;
            }
            .lang-selectors span {
                margin: 10px 0;
            }
            .input-group {
                flex-direction: column;
                gap: 10px;
            }
            .mic-btn {
                width: 48px;
                height: 48px;
                min-width: 48px;
                font-size: 24px;
                align-self: center;
            }
            textarea {
                min-height: 100px;
                font-size: 16px;
            }
            button[type="submit"], .action-button {
                padding: 14px 0;
                font-size: 18px;
            }
            .result {
                padding: 18px 15px;
                font-size: 18px;
            }
            .image-display-section {
                flex-direction: column;
                gap: 18px;
            }
            .image-container {
                max-width: 100%;
            }
            .image-translator-section input[type="file"],
            .image-translator-section .action-button {
                width: 100%;
                max-width: none;
            }
            .file-translator-section #download-link, .file-translator-section #view-button {
                margin-right: 0;
                width: 100%;
            }
        }

        @media (max-width: 480px) {
            body {
                padding: 10px;
            }
            .container {
                padding: 20px 10px;
            }
            h1 {
                font-size: 1.8rem;
            }
            .tab-button {
                font-size: 0.9em;
            }
            .lang-selectors select {
                font-size: 15px;
            }
            textarea {
                font-size: 15px;
            }
            button[type="submit"], .action-button {
                font-size: 16px;
            }
            .result {
                font-size: 16px;
            }
        }
        /* --- Hiệu ứng hover & chuyển động cho Dịch Hình Ảnh và Dịch File --- */

/* Nút hành động (ảnh & file) */
.image-translator-section .action-button,
#startWebcamButton, #captureButton,
.file-translator-section .action-button,
#download-link, #view-button {
    transition: 
        background 0.3s, 
        color 0.3s, 
        box-shadow 0.3s, 
        transform 0.2s;
    box-shadow: 0 2px 10px rgba(56,198,229,0.10);
    position: relative;
    overflow: hidden;
}
.image-translator-section .action-button:hover,
#startWebcamButton:hover, #captureButton:hover,
.file-translator-section .action-button:hover,
#download-link:hover, #view-button:hover {
    background: linear-gradient(100deg, #38c6e5 0%, #5b6ee1 100%);
    color: #fff;
    box-shadow: 0 8px 24px rgba(91,110,225,0.18);
    transform: translateY(-2px) scale(1.04);
}
.image-translator-section .action-button:active,
#startWebcamButton:active, #captureButton:active,
.file-translator-section .action-button:active,
#download-link:active, #view-button:active {
    transform: scale(0.97);
}

/* Khung ảnh */
.image-container {
    transition: box-shadow 0.3s, transform 0.2s, border-color 0.3s;
}
.image-container:hover {
    box-shadow: 0 8px 32px rgba(56,198,229,0.18), 0 2px 8px rgba(91,110,225,0.13);
    transform: translateY(-4px) scale(1.03) rotate(-1deg);
    border-color: #38c6e5;
}
.image-container img, .image-container video {
    transition: box-shadow 0.3s, filter 0.3s;
}
.image-container img:hover, .image-container video:hover {
    box-shadow: 0 8px 32px rgba(91,110,225,0.18);
    filter: brightness(1.08) saturate(1.2);
}

/* Kết quả dịch văn bản từ ảnh */
#translatedTextResults {
    animation: fadeIn 1s;
    border-left: 6px solid #38c6e5;
    box-shadow: 0 4px 18px rgba(56,198,229,0.10);
    transition: box-shadow 0.3s, border-color 0.3s;
}
#translatedTextResults:hover {
    border-color: #5b6ee1;
    box-shadow: 0 8px 32px rgba(91,110,225,0.13);
}
#resultsList li {
    transition: background 0.3s, border-left 0.3s, transform 0.2s;
    cursor: pointer;
}
#resultsList li:hover {
    background: #b3e5fc;
    border-left: 8px solid #5b6ee1;
    transform: scale(1.03) translateX(4px);
}

/* File translator section */
.file-translator-section input[type="file"]:hover {
    border-color: #38c6e5;
    background: #eaf6ff;
}
.file-translator-section select:focus {
    border-color: #5b6ee1;
    background: #eaf6ff;
    box-shadow: 0 0 0 3px #38c6e522;
}
.file-translator-section #translated-text-display {
    animation: fadeIn 1s;
    border-left: 6px solid #38c6e5;
    box-shadow: 0 4px 18px rgba(56,198,229,0.10);
    transition: box-shadow 0.3s, border-color 0.3s;
}
.file-translator-section #translated-text-display:hover {
    border-color: #5b6ee1;
    box-shadow: 0 8px 32px rgba(91,110,225,0.13);
}

/* Hiệu ứng xuất hiện */
@keyframes fadeIn {
    from { opacity: 0; transform: translateY(20px);}
    to { opacity: 1; transform: translateY(0);}
}

/* --- Đẹp hóa input file, select, input number cho dịch ảnh & file --- */

/* Input file đẹp */
.image-translator-section input[type="file"],
.file-translator-section input[type="file"] {
    display: block;
    margin: 0 auto 18px auto;
    padding: 12px 16px;
    border: 2px dashed var(--secondary-color);
    border-radius: 10px;
    background: #fafdff;
    font-size: 1em;
    width: 80%;
    max-width: 400px;
    color: var(--primary-color);
    cursor: pointer;
    transition: 
        border-color 0.3s, 
        background 0.3s, 
        box-shadow 0.3s;
    box-shadow: 0 2px 10px rgba(56,198,229,0.06);
}
.image-translator-section input[type="file"]:hover,
.file-translator-section input[type="file"]:hover {
    border-color: #5b6ee1;
    background: #eaf6ff;
    box-shadow: 0 4px 18px rgba(91,110,225,0.13);
}
.image-translator-section input[type="file"]:focus,
.file-translator-section input[type="file"]:focus {
    outline: none;
    border-color: #38c6e5;
    background: #eaf6ff;
    box-shadow: 0 0 0 3px #38c6e522;
}

/* Select đẹp */
.image-translator-section select,
.file-translator-section select {
    width: 80%;
    max-width: 400px;
    padding: 12px 16px;
    border: 2px solid var(--border-color);
    border-radius: 10px;
    font-size: 1em;
    background: #fafdff;
    color: var(--primary-color);
    margin-bottom: 12px;
    transition: 
        border-color 0.3s, 
        background 0.3s, 
        box-shadow 0.3s;
    box-shadow: 0 1px 6px rgba(91,110,225,0.04);
    appearance: none;
    -webkit-appearance: none;
    -moz-appearance: none;
    background-image: url("data:image/svg+xml,%3csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 16 16'%3e%3cpath fill='none' stroke='%235b6ee1' stroke-linecap='round' stroke-linejoin='round' stroke-width='2' d='M2 5l6 6 6-6'/%3e%3c/svg%3e");
    background-repeat: no-repeat;
    background-position: right 1rem center;
    background-size: 18px 14px;
}
.image-translator-section select:focus,
.file-translator-section select:focus {
    border-color: #38c6e5;
    background: #eaf6ff;
    box-shadow: 0 0 0 3px #38c6e522;
    outline: none;
}
.image-translator-section select:hover,
.file-translator-section select:hover {
    background: #f0fcfa;
    border-color: #5b6ee1;
}

/* Input chỉnh size chữ */
.image-translator-section input[type="number"] {
    width: 80%;
    max-width: 400px;
    padding: 12px 16px;
    border: 2px solid var(--border-color);
    border-radius: 10px;
    font-size: 1em;
    background: #fafdff;
    color: var(--primary-color);
    margin-bottom: 12px;
    transition: 
        border-color 0.3s, 
        background 0.3s, 
        box-shadow 0.3s;
    box-shadow: 0 1px 6px rgba(91,110,225,0.04);
}
.image-translator-section input[type="number"]:focus {
    border-color: #38c6e5;
    background: #eaf6ff;
    box-shadow: 0 0 0 3px #38c6e522;
    outline: none;
}
.image-translator-section input[type="number"]:hover {
    background: #f0fcfa;
    border-color: #5b6ee1;
}

/* Label căn trái, đậm và đồng bộ màu */
.image-translator-section label,
.file-translator-section label {
    display: block;
    font-weight: 600;
    color: var(--primary-color);
    margin-bottom: 6px;
    margin-left: 8px;
    text-align: left;
    letter-spacing: 0.2px;
}
    </style>
</head>
<body>
    <div class="container">
        <h1>Dịch thuật Đa Năng 🌐</h1>

        <div class="tab-container">
            <button class="tab-button active" data-tab="text-translator">Dịch Văn Bản</button>
            <button class="tab-button" data-tab="image-translator">Dịch Hình Ảnh</button>
            <button class="tab-button" data-tab="file-translator">Dịch File</button>
        </div>

        <!-- Tab Dịch Văn Bản -->
        <div id="text-translator" class="tab-content active">
            <form id="translate-text-form">
                <div class="lang-selectors">
                    <select id="source-lang">
                        <option value="ja">Nhật</option>
                        <option value="vi">Việt</option>
                        <option value="en">Anh</option>
                        <option value="ko">Hàn</option>
                        <option value="zh-cn">Trung (Giản thể)</option>
                        <option value="fr">Pháp</option>
                        <option value="de">Đức</option>
                        <option value="es">Tây Ban Nha</option>
                    </select>
                    <button type="button" id="swap-btn" class="swap-btn" title="Đổi chiều">
                        <span class="swap-icon">&#8646;</span>
                    </button>
                    <select id="dest-lang">
                        <option value="vi" selected>Việt</option>
                        <option value="ja">Nhật</option>
                        <option value="en">Anh</option>
                        <option value="ko">Hàn</option>
                        <option value="zh-cn">Trung (Giản thể)</option>
                        <option value="fr">Pháp</option>
                        <option value="de">Đức</option>
                        <option value="es">Tây Ban Nha</option>
                    </select>
                </div>
                <div class="input-group">
                    <textarea id="source-text" placeholder="Nhập văn bản..."></textarea>
                    <button type="button" class="mic-btn" id="mic-btn" title="Nói"><span id="mic-icon">🎤</span></button>
                </div>
                <button type="submit" class="action-button">Dịch</button>
            </form>
            <div class="result" id="text-translation-result"></div>
        </div>

        <!-- Tab Dịch Hình Ảnh -->
        <div id="image-translator" class="tab-content">
            <p>Tải lên một hình ảnh hoặc sử dụng webcam để phát hiện và dịch văn bản trong đó.</p>

            <div class="input-group">
                <input type="file" id="imageUpload" accept="image/*">
                <button id="translateImageButton" class="action-button" disabled>Dịch Hình Ảnh Đã Tải Lên</button>
            </div>

            <hr>

            <div class="input-group">
                <button id="startWebcamButton" class="action-button">Bật Webcam</button>
                <button id="captureButton" class="action-button" disabled>Chụp và Dịch</button>
            </div>
            <video id="webcamFeed" autoplay playsinline class="hidden"></video>
            <canvas id="canvas" class="hidden"></canvas>

            <div class="input-group">
                <label for="imageLanguageSelect">Ngôn ngữ đích:</label>
                <select id="imageLanguageSelect">
                    <option value="vi">Tiếng Việt (vi)</option>
                    <option value="en">Tiếng Anh (en)</option>
                    <option value="ja">Tiếng Nhật (ja)</option>
                    <option value="fr">Tiếng Pháp (fr)</option>
                    <option value="de">Tiếng Đức (de)</option>
                    <option value="es">Tiếng Tây Ban Nha (es)</option>
                </select>
            </div>
            <div class="input-group">
                <label for="fontSizeInput">Cỡ chữ:</label>
                <input type="number" id="fontSizeInput" value="18" min="8" max="72">
            </div>

            <div id="loading-image" class="hidden">Đang xử lý...</div>

            <div class="image-display-section">
                <div class="image-container">
                    <h2>Xem Trước Ảnh Gốc</h2>
                    <img id="originalImagePreview" src="#" alt="Original Image" class="hidden">
                </div>
                <div class="image-container">
                    <h2>Ảnh Đã Dịch</h2>
                    <img id="translatedImageDisplay" src="#" alt="Translated Image" class="hidden">
                </div>
            </div>

            <div id="translatedTextResults" class="hidden">
                <h2>Kết Quả Dịch Văn Bản</h2>
                <ul id="resultsList"></ul>
            </div>
        </div>

        <!-- Tab Dịch File -->
        <div id="file-translator" class="tab-content">
            <p>Tải lên file của bạn, chúng tôi sẽ dịch và giữ nguyên định dạng nhiều nhất có thể.</p>
            <form id="translate-file-form">
                <div class="form-group">
                    <label for="file-upload">Chọn File:</label>
                    <input type="file" id="file-upload" name="file" accept=".txt,.docx,.pptx,.pdf" required>
                </div>
                <div class="form-group">
                    <label for="file-language-select">Dịch sang ngôn ngữ:</label>
                    <select id="file-language-select" name="language">
                        <option value="en">Tiếng Anh</option>
                        <option value="vi">Tiếng Việt</option>
                        <option value="fr">Tiếng Pháp</option>
                        <option value="es">Tiếng Tây Ban Nha</option>
                        <option value="ja">Tiếng Nhật</option>
                        <option value="ko">Tiếng Hàn</option>
                        <option value="zh-CN">Tiếng Trung (Giản thể)</option>
                        <option value="de">Tiếng Đức</option>
                    </select>
                </div>
                <button type="submit" class="action-button">Dịch File</button>
            </form>
            <div id="loading-file" class="hidden">
                <div class="spinner"></div>
                <p id="status-file">Đang xử lý, vui lòng chờ...</p>
            </div>
            <div id="translated-text-display" class="hidden"></div>
            <div id="download-section" class="hidden">
                <p>File của bạn đã được dịch thành công!</p>
                <a id="view-button" href="#" target="_blank">Xem File</a>
                <a id="download-link" href="#">Tải về</a>
            </div>
        </div>
    </div>

    <script>
    document.addEventListener('DOMContentLoaded', () => {
            // --- Quản lý Tabs ---
            const tabButtons = document.querySelectorAll('.tab-button');
            const tabContents = document.querySelectorAll('.tab-content');

            tabButtons.forEach(button => {
                button.addEventListener('click', () => {
                    tabButtons.forEach(btn => btn.classList.remove('active'));
                    tabContents.forEach(content => content.classList.remove('active'));

                    button.classList.add('active');
                    document.getElementById(button.dataset.tab).classList.add('active');
                });
            });

            // --- Logic cho Dịch Văn Bản (index.py) ---
            const sourceLangSelect = document.getElementById('source-lang');
            const destLangSelect = document.getElementById('dest-lang');
            const textTranslateForm = document.getElementById('translate-text-form');
            const textarea = document.getElementById('source-text');
            const textResultDiv = document.getElementById('text-translation-result');
            const micBtn = document.getElementById('mic-btn');
            const micIcon = document.getElementById('mic-icon');
            const swapBtn = document.getElementById('swap-btn');

            async function performTextTranslation(textToTranslate) {
                if (!textToTranslate.trim()) {
                    textResultDiv.innerText = 'Vui lòng nhập văn bản.';
                    return null;
                }
                textResultDiv.innerText = 'Đang dịch...';
                const srcLang = sourceLangSelect.value;
                const destLang = destLangSelect.value;

                try {
                    const res = await fetch('/translate_text', {
                        method: 'POST',
                        headers: { 'Content-Type': 'application/json' },
                        body: JSON.stringify({ text: textToTranslate, src: srcLang, dest: destLang })
                    });
                    if (!res.ok) throw new Error(`Lỗi server: ${res.status}`);
                    const data = await res.json();
                    const translatedText = data.translated || 'Không thể dịch.';
                    textResultDiv.innerText = translatedText;
                    return translatedText;
                } catch (error) {
                    console.error('Translation failed:', error);
                    textResultDiv.innerText = 'Đã xảy ra lỗi khi dịch.';
                    return null;
                }
            }

            textTranslateForm.addEventListener('submit', (e) => {
                e.preventDefault();
                performTextTranslation(textarea.value);
            });

            let voices = [];
            function populateVoiceList() {
                voices = window.speechSynthesis.getVoices();
            }
            populateVoiceList();
            if (speechSynthesis.onvoiceschanged !== undefined) {
                speechSynthesis.onvoiceschanged = populateVoiceList;
            }

            function speak(text, lang) {
                if (!text || !window.speechSynthesis) return;
                const utterance = new SpeechSynthesisUtterance(text);
                const targetVoice = voices.find(voice => voice.lang.startsWith(lang));
                utterance.lang = lang;
                if (targetVoice) utterance.voice = targetVoice;
                window.speechSynthesis.speak(utterance);
            }

            const speechLangMap = { 'ja': 'ja-JP', 'vi': 'vi-VN', 'en': 'en-US', 'ko': 'ko-KR', 'zh-cn': 'cmn-Hans-CN', 'fr': 'fr-FR', 'de': 'de-DE', 'es': 'es-ES' };
            let recognizing = false;
            let recognition;

            if ('webkitSpeechRecognition' in window) {
                recognition = new webkitSpeechRecognition();
                recognition.continuous = false;
                recognition.interimResults = true; // Đã thay đổi: Bật kết quả tạm thời

                let finalTranscript = ''; // Biến để lưu trữ bản ghi cuối cùng

                recognition.onstart = () => {
                    recognizing = true;
                    micBtn.classList.add('listening');
                    micIcon.textContent = '🔴';
                    finalTranscript = ''; // Reset bản ghi cuối cùng khi bắt đầu
                    textarea.value = ''; // Xóa textarea khi bắt đầu
                };
                recognition.onend = () => {
                    recognizing = false;
                    micBtn.classList.remove('listening');
                    micIcon.textContent = '🎤';
                    // Nếu không có kết quả cuối cùng nào được ghi lại (ví dụ, người dùng dừng nói đột ngột)
                    if (finalTranscript.trim() === '' && textarea.value.trim() !== '') {
                        // Nếu textarea đã có nội dung từ interim results nhưng không có final result,
                        // thì coi nội dung hiện tại trong textarea là bản ghi cuối cùng.
                        performTextTranslation(textarea.value);
                    } else if (finalTranscript.trim() !== '') {
                        // Nếu có bản ghi cuối cùng, tiến hành dịch
                        performTextTranslation(finalTranscript);
                    }
                };
                recognition.onerror = (event) => {
                    console.error('Speech recognition error:', event.error);
                    micBtn.classList.remove('listening');
                    micIcon.textContent = '🎤';
                    alert('Lỗi nhận dạng giọng nói: ' + event.error);
                };

                // Đã thêm: Xử lý kết quả tạm thời
                recognition.onresult = async (event) => {
                    let interimTranscript = '';
                    for (let i = event.resultIndex; i < event.results.length; ++i) {
                        if (event.results[i].isFinal) {
                            finalTranscript += event.results[i][0].transcript;
                        } else {
                            interimTranscript += event.results[i][0].transcript;
                        }
                    }
                    // Hiển thị cả bản ghi cuối cùng và bản ghi tạm thời trong textarea
                    textarea.value = finalTranscript + interimTranscript;

                    // Nếu đây là kết quả cuối cùng, tiến hành dịch và phát âm
                    if (event.results[event.resultIndex].isFinal) {
                        const translatedText = await performTextTranslation(finalTranscript);
                        if (translatedText && !translatedText.includes('lỗi') && !translatedText.includes('Không thể')) {
                            speak(translatedText, destLangSelect.value);
                        }
                    }
                };

                micBtn.onclick = () => {
                    if (recognizing) {
                        recognition.stop();
                    } else {
                        recognition.lang = speechLangMap[sourceLangSelect.value] || 'en-US';
                        try {
                            recognition.start();
                        } catch (e) {
                            alert("Không thể bắt đầu nhận dạng giọng nói. Vui lòng kiểm tra quyền truy cập microphone.");
                            console.error("Lỗi khi bắt đầu nhận dạng giọng nói:", e);
                        }
                    }
                };
            } else {
                micBtn.disabled = true;
                micBtn.title = 'Trình duyệt không hỗ trợ nhận diện giọng nói';
                micIcon.textContent = '🚫';
            }

            function updateUIForLangChange() {
                const langName = sourceLangSelect.options[sourceLangSelect.selectedIndex].text;
                textarea.placeholder = `Nhập hoặc nói tiếng ${langName}...`;
                micBtn.title = `Nói tiếng ${langName}`;
            }
            sourceLangSelect.addEventListener('change', updateUIForLangChange);

            swapBtn.onclick = function () {
                swapBtn.classList.add('swapping');
                sourceLangSelect.classList.add('swapping');
                destLangSelect.classList.add('swapping');
                setTimeout(() => {
                    swapBtn.classList.remove('swapping');
                    sourceLangSelect.classList.remove('swapping');
                    destLangSelect.classList.remove('swapping');
                    const temp = sourceLangSelect.value;
                    sourceLangSelect.value = destLangSelect.value;
                    destLangSelect.value = temp;
                    updateUIForLangChange();
                    textarea.value = '';
                    textResultDiv.innerText = '';
                    sourceLangSelect.classList.add('swapped');
                    destLangSelect.classList.add('swapped');
                    setTimeout(() => {
                        sourceLangSelect.classList.remove('swapped');
                        destLangSelect.classList.remove('swapped');
                    }, 600);
                }, 500);
            };

            // --- Logic cho Dịch Hình Ảnh (main.py) ---
            const imageUpload = document.getElementById('imageUpload');
            const imageLanguageSelect = document.getElementById('imageLanguageSelect');
            const fontSizeInput = document.getElementById('fontSizeInput');
            const translateImageButton = document.getElementById('translateImageButton');
            const startWebcamButton = document.getElementById('startWebcamButton');
            const captureButton = document.getElementById('captureButton');
            const webcamFeed = document.getElementById('webcamFeed');
            const canvas = document.getElementById('canvas');
            const context = canvas.getContext('2d');

            const originalImagePreview = document.getElementById('originalImagePreview');
            const translatedImageDisplay = document.getElementById('translatedImageDisplay');
            const loadingImageDiv = document.getElementById('loading-image');
            const translatedTextResultsDiv = document.getElementById('translatedTextResults');
            const resultsList = document.getElementById('resultsList');

            let selectedFile = null;
            let mediaStream = null; // Để giữ luồng webcam

            // Hàm để đặt lại hiển thị hình ảnh
            const resetImageDisplays = () => {
                originalImagePreview.classList.add('hidden');
                originalImagePreview.src = '#';
                translatedImageDisplay.classList.add('hidden');
                translatedImageDisplay.src = '#';
                translatedTextResultsDiv.classList.add('hidden');
                resultsList.innerHTML = '';
            };

            // Lắng nghe sự kiện tải file hình ảnh
            imageUpload.addEventListener('change', (event) => {
                selectedFile = event.target.files[0];
                if (selectedFile) {
                    resetImageDisplays(); // Xóa kết quả trước đó
                    const reader = new FileReader();
                    reader.onload = (e) => {
                        originalImagePreview.src = e.target.result;
                        originalImagePreview.classList.remove('hidden');
                        translateImageButton.disabled = false;
                        captureButton.disabled = true; // Vô hiệu hóa nút chụp nếu đang tải lên
                        stopWebcam(); // Dừng webcam nếu đang hoạt động
                    };
                    reader.readAsDataURL(selectedFile);
                } else {
                    translateImageButton.disabled = true;
                    resetImageDisplays();
                }
            });

            // Lắng nghe sự kiện nút "Dịch Hình Ảnh Đã Tải Lên"
            translateImageButton.addEventListener('click', async () => {
                if (!selectedFile) {
                    alert('Vui lòng chọn một hình ảnh trước.');
                    return;
                }
                await processImage(selectedFile);
            });

            // Lắng nghe sự kiện nút "Bật Webcam"
            startWebcamButton.addEventListener('click', async () => {
                resetImageDisplays(); // Xóa kết quả trước đó
                selectedFile = null; // Xóa bất kỳ file nào đã chọn trước đó
                translateImageButton.disabled = true; // Vô hiệu hóa nút dịch tải lên

                if (mediaStream) {
                    stopWebcam(); // Dừng nếu đang chạy
                    startWebcamButton.textContent = 'Bật Webcam';
                    captureButton.disabled = true;
                    webcamFeed.classList.add('hidden');
                    return;
                }

                try {
                    mediaStream = await navigator.mediaDevices.getUserMedia({ video: true });
                    webcamFeed.srcObject = mediaStream;
                    webcamFeed.classList.remove('hidden');
                    captureButton.disabled = false;
                    startWebcamButton.textContent = 'Dừng Webcam';
                    originalImagePreview.classList.add('hidden'); // Ẩn xem trước ảnh đã tải lên
                } catch (err) {
                    console.error('Lỗi truy cập webcam:', err);
                    alert('Không thể truy cập webcam. Vui lòng đảm bảo bạn có camera và đã cấp quyền.');
                    captureButton.disabled = true;
                    webcamFeed.classList.add('hidden');
                    startWebcamButton.textContent = 'Bật Webcam';
                }
            });

            // Hàm để dừng luồng webcam
            const stopWebcam = () => {
                if (mediaStream) {
                    mediaStream.getTracks().forEach(track => track.stop());
                    webcamFeed.srcObject = null;
                    mediaStream = null;
                }
            };

            // Lắng nghe sự kiện nút "Chụp và Dịch"
            captureButton.addEventListener('click', async () => {
                if (!mediaStream) {
                    alert('Webcam không hoạt động. Vui lòng bật webcam trước.');
                    return;
                }

                // Đặt kích thước canvas khớp với luồng video
                canvas.width = webcamFeed.videoWidth;
                canvas.height = webcamFeed.videoHeight;
                context.drawImage(webcamFeed, 0, 0, canvas.width, canvas.height);

                // Lấy dữ liệu hình ảnh từ canvas dưới dạng Blob
                canvas.toBlob(async (blob) => {
                    if (blob) {
                        originalImagePreview.src = URL.createObjectURL(blob); // Hiển thị ảnh đã chụp trong xem trước
                        originalImagePreview.classList.remove('hidden');
                        await processImage(blob); // Xử lý blob đã chụp
                    } else {
                        alert('Không thể chụp ảnh từ webcam.');
                    }
                }, 'image/png');
            });

            // Hàm chung để xử lý và dịch hình ảnh (tải lên hoặc chụp)
            async function processImage(imageBlobOrFile) {
                loadingImageDiv.classList.remove('hidden');
                translateImageButton.disabled = true;
                captureButton.disabled = true;
                startWebcamButton.disabled = true; // Vô hiệu hóa nút webcam trong quá trình xử lý

                const formData = new FormData();
                formData.append('image', imageBlobOrFile);
                formData.append('language', imageLanguageSelect.value);
                formData.append('fontSize', fontSizeInput.value);

                try {
                    const response = await fetch('/translate_image', {
                        method: 'POST',
                        body: formData
                    });

                    if (!response.ok) {
                        const errorData = await response.json();
                        throw new Error(errorData.error || 'Phản hồi mạng không ổn định');
                    }

                    const data = await response.json();

                    if (data.processed_image) {
                        translatedImageDisplay.src = `data:image/png;base64,${data.processed_image}`;
                        translatedImageDisplay.classList.remove('hidden');
                    } else {
                        translatedImageDisplay.classList.add('hidden');
                    }

                    if (data.detections && data.detections.length > 0) {
                        resultsList.innerHTML = '';
                        data.detections.forEach(detection => {
                            const listItem = document.createElement('li');
                            listItem.innerHTML = `<strong>Gốc:</strong> ${detection.original}<br><strong>Dịch:</strong> ${detection.translated}<br><em>Độ tin cậy: ${detection.confidence}</em>`;
                            resultsList.appendChild(listItem);
                        });
                        translatedTextResultsDiv.classList.remove('hidden');
                    } else {
                        resultsList.innerHTML = '<li>Không phát hiện hoặc dịch được văn bản với ngưỡng tin cậy đã cho.</li>';
                        translatedTextResultsDiv.classList.remove('hidden');
                    }

                } catch (error) {
                    console.error('Lỗi:', error);
                    alert(`Lỗi khi dịch hình ảnh: ${error.message}`);
                    translatedImageDisplay.classList.add('hidden');
                    translatedTextResultsDiv.classList.add('hidden');
                } finally {
                    loadingImageDiv.classList.add('hidden');
                    // Bật lại các nút dựa trên trạng thái hiện tại (webcam hoặc tải file)
                    if (selectedFile) { // Nếu một hình ảnh đã được tải lên
                        translateImageButton.disabled = false;
                        captureButton.disabled = true;
                        startWebcamButton.disabled = false;
                    } else if (mediaStream) { // Nếu webcam đang hoạt động
                        translateImageButton.disabled = true;
                        captureButton.disabled = false;
                        startWebcamButton.disabled = false;
                    } else { // Không có hình ảnh hoặc webcam hoạt động
                        translateImageButton.disabled = true;
                        captureButton.disabled = true;
                        startWebcamButton.disabled = false;
                    }
                }
            }

            // --- Logic cho Dịch File (fiel.py) ---
            const fileTranslateForm = document.getElementById('translate-file-form');
            const fileInput = document.getElementById('file-upload');
            const fileLanguageSelect = document.getElementById('file-language-select');
            const loadingFileElem = document.getElementById('loading-file');
            const statusFileElem = document.getElementById('status-file');
            const submitFileButton = fileTranslateForm.querySelector('button[type="submit"]');
            const downloadSection = document.getElementById('download-section');
            const downloadLink = document.getElementById('download-link');
            const viewButton = document.getElementById('view-button');
            const translatedTextDisplay = document.getElementById('translated-text-display');

            fileTranslateForm.addEventListener('submit', async function(event) {
                event.preventDefault();

                // Reset trạng thái trước đó
                downloadSection.classList.add('hidden');
                translatedTextDisplay.classList.add('hidden');
                translatedTextDisplay.innerHTML = ''; // Xóa văn bản trước đó
                downloadLink.removeAttribute('href');
                viewButton.removeAttribute('href');

                if (fileInput.files.length === 0) {
                    alert('Vui lòng chọn một file để dịch.');
                    return;
                }

                const originalFileName = fileInput.files[0].name;
                const fileExtension = originalFileName.split('.').pop().toLowerCase();

                loadingFileElem.classList.remove('hidden');
                submitFileButton.disabled = true;
                submitFileButton.style.opacity = '0.6';
                statusFileElem.textContent = 'Đang tải file lên...';

                const formData = new FormData();
                formData.append('file', fileInput.files[0]);
                formData.append('language', fileLanguageSelect.value);

                try {
                    statusFileElem.textContent = 'File đang được dịch, quá trình này có thể mất vài phút...';
                    const response = await fetch('/translate_file', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = window.URL.createObjectURL(blob);

                        let mimeType;
                        // Đặt loại MIME phù hợp để xem
                        if (fileExtension === 'pdf') {
                            mimeType = 'application/pdf';
                        } else if (fileExtension === 'txt') {
                            mimeType = 'text/plain';
                            // Đối với file TXT, đọc nội dung blob để hiển thị
                            const textContent = await blob.text();
                            translatedTextDisplay.textContent = textContent;
                            translatedTextDisplay.classList.remove('hidden');
                        } else if (fileExtension === 'docx') {
                            mimeType = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document';
                        } else if (fileExtension === 'pptx') {
                            mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation';
                        } else {
                            mimeType = 'application/octet-stream'; // Dự phòng
                        }

                        // Tạo một blob với loại MIME chính xác để xem
                        const viewBlob = new Blob([blob], { type: mimeType });
                        const viewUrl = window.URL.createObjectURL(viewBlob);

                        downloadLink.href = url;
                        downloadLink.download = 'translated_' + originalFileName;
                        viewButton.href = viewUrl;

                        statusFileElem.textContent = 'Dịch thành công! Bạn có thể xem hoặc tải file.';
                        downloadSection.classList.remove('hidden');

                    } else {
                        const errorResult = await response.json();
                        alert('Lỗi: ' + errorResult.error);
                        statusFileElem.textContent = 'Đã xảy ra lỗi.';
                    }
                } catch (error) {
                    alert('Lỗi kết nối hoặc xử lý file. Vui lòng thử lại.');
                    console.error('Error:', error);
                    statusFileElem.textContent = 'Đã xảy ra lỗi.';
                } finally {
                    loadingFileElem.classList.add('hidden');
                    submitFileButton.disabled = false;
                    submitFileButton.style.opacity = '1';
                }
            });
        });
    </script>
</body>
</html>
"""

@app.route('/')
def index():
    return render_template_string(HTML_TEMPLATE)


@app.route("/translate_text", methods=["POST"]) 
def translate_text_route():
    data = request.get_json()
    if not data:
        return jsonify({"translated": "Lỗi: Yêu cầu không hợp lệ."}), 400

    text = data.get("text", "").strip()
    src_lang = data.get("src", "ja")
    dest_lang = data.get("dest", "vi")

    if not text:
        return jsonify({"translated": ""})

    cache_key = f"{src_lang}:{dest_lang}:{text.lower()}"

    if cache_key in translation_cache:
        translation_cache.move_to_end(cache_key)
        return jsonify({"translated": translation_cache[cache_key]})

    try:
        result = text_translator.translate(text, src=src_lang, dest=dest_lang)
        translated_text = result.text

        translation_cache[cache_key] = translated_text
        if len(translation_cache) > MAX_CACHE_SIZE:
            translation_cache.popitem(last=False)

        return jsonify({"translated": translated_text})
    except Exception as e:
        print(f"Error during translation: {e}")
        return jsonify({"translated": "Lỗi: Không thể dịch văn bản."}), 500


@app.route('/translate_image', methods=['POST'])
def translate_image_route():
    if 'image' not in request.files:
        return jsonify({'error': 'Không có file hình ảnh được cung cấp'}), 400

    file = request.files['image']
    target_language = request.form.get('language', 'vi')
    font_size = int(request.form.get('fontSize', 18))

    try:
        img_bytes = file.read()
        np_arr = np.frombuffer(img_bytes, np.uint8)
        img = cv2.imdecode(np_arr, cv2.IMREAD_COLOR)

        if img is None:
            return jsonify({'error': 'Không thể giải mã hình ảnh'}), 400

        img_translated = img.copy()

        detections = reader.readtext(img)

        translated_texts_info = []

        for i, (bbox_points, text, score) in enumerate(detections):
            if score > 0.25:
                bbox = np.array(bbox_points, dtype=np.int32)
                translated_text = translate_text_image(text, target_language)

                translated_texts_info.append({
                    'original': text,
                    'translated': translated_text,
                    'confidence': f"{score:.2f}"
                })

                img_translated = blur_bbox(img_translated, bbox)

                top_left_corner = tuple(bbox[0])

                img_translated = draw_text_utf8(img_translated, translated_text, top_left_corner, FONT_PATH, font_size, color=(255, 255, 255))

        is_success, buffer = cv2.imencode(".png", img_translated)
        if not is_success:
            return jsonify({'error': 'Không thể mã hóa hình ảnh đã xử lý'}), 500

        img_base64 = base64.b64encode(buffer).decode('utf-8')

        return jsonify({
            'processed_image': img_base64,
            'detections': translated_texts_info
        })

    except Exception as e:
        print(f"Lỗi xử lý hình ảnh: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/translate_file', methods=['POST']) 
def translate_file_route():
    if 'file' not in request.files:
        return jsonify({'error': 'Không tìm thấy file'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'Chưa chọn file nào'}), 400

    target_lang = request.form.get('language', 'en')
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()

    try:
        file_stream = BytesIO(file.read())
        output_stream = None

        if file_ext == '.txt':
            output_stream = handle_txt(file_stream, target_lang)
        elif file_ext == '.docx':
            output_stream = handle_docx(file_stream, target_lang)
        elif file_ext == '.pptx':
            output_stream = handle_pptx(file_stream, target_lang)
        else:
            return jsonify({'error': f'Định dạng file {file_ext} không được hỗ trợ'}), 415

        if file_ext == '.docx':
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif file_ext == '.pptx':
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif file_ext == '.txt':
            mime_type = 'text/plain'
        else:
            mime_type = 'application/octet-stream'

        return send_file(
            output_stream,
            mimetype=mime_type,
        )

    except Exception as e:
        print(f"Lỗi xử lý file: {e}")
        return jsonify({'error': f'Đã xảy ra lỗi khi xử lý file: {str(e)}'}), 500

if __name__ == '__main__':
    print("Khởi động server...")
    print("LƯU Ý QUAN TRỌNG CHO VIỆC DỊCH PDF:")
    print("1. Tải font 'DejaVuSans.ttf' từ https://dejavu-fonts.github.io/")
    print("2. Đặt file font đó vào cùng thư mục với file script này.")
    print("LƯU Ý QUAN TRỌNG CHO VIỆC DỊCH HÌNH ẢNH:")
    print("1. Tải font 'NotoSans-Regular.ttf' từ https://fonts.google.com/noto/specimen/Noto+Sans")
    print("2. Đặt file font đó vào cùng thư mục với file script này.")
import os  # Thêm dòng này ở đầu nếu chưa có

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)  # Tắt debug cho production